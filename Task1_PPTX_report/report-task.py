import sys
import json
import csv
import matplotlib.pyplot as plt
# pptx throws attribute error on python 3.10 (on windows), necessary imports below
import collections
import collections.abc
c = collections
c.abc = collections.abc
# # # #
from pptx import Presentation
from pptx.util import Inches


config_file = sys.argv[1]

with open (config_file, "r") as f:
    data = json.load(f)
    prs = Presentation()

    def title_slide(data):
        # title-subtitle layout
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data["title"]
        subtitle.text = data["content"]
        
    def text_slide(data):
        # Title only layout
        text_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(text_slide_layout)
        title = slide.shapes.title
        title.text = data["title"]
        # Add textbox separately
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = data["content"]

    def bullet_point_slide(data):
        # Title - bullet point layout
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title = shapes.title
        body = shapes.placeholders[1]
        title.text = data["title"]
        try:
            for levels in data["content"]:
                tf = body.text_frame
                if(levels["level"] == 1):
                    tf = body.text_frame
                    p = tf.add_paragraph()
                    p.text = levels["text"]
                    p.level = levels["level"] 
                    
                elif(levels["level"] == 2):
                    tf = body.text_frame
                    p = tf.add_paragraph()
                    p.text = levels["text"]
                    p.level = levels["level"]
                    
                elif(levels["level"] == 3):
                    tf = body.text_frame
                    p = tf.add_paragraph()
                    p.text = levels["text"]
                    p.level = levels["level"]
        except:
            print("Level amount out of bound")
            
    def picture_slide(data):
        # Title only layout
        picture_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(picture_slide_layout)
        title = slide.shapes.title
        title.text = data["title"]
        try:
            img_path = data["content"]
            # Change left for indentation, change top for distance from top, change height for obviously, the height
            left = top = Inches(1.3)
            height = Inches(4)
            slide.shapes.add_picture(img_path, left, top, height=height)
        except:
            
            print("No matching picture found in folder")
    def plot_slide(data):
        plot_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(plot_slide_layout)
        title = slide.shapes.title
        title.text = data["title"]
        data_file = data["content"]
        # Convert content name of the object to csv
        try:
            csv_file = data_file.replace(".dat",".csv")
            x_values = []
            y_values = []
            #open 
            with open(csv_file, "r") as file:
                reader = csv.reader(file, delimiter=";")
                for row in reader:
                    x_values.append(float(row[0]))
                    y_values.append(float(row[1]))

            plt.plot(x_values, y_values)
            plt.xlabel(data["configuration"]["x-label"])
            plt.ylabel(data["configuration"]["y-label"])
            # For every plot, there will be a png saved to the folder, unsure whether they are wanted or not
            plt.savefig("temp_plot.png")
            img_path = "temp_plot.png"
            left = top = Inches(1.3)
            height = Inches(4.5)
            slide.shapes.add_picture(img_path, left, top, height=height)
        except:
            print("No dat file found")

    # Loop through the json, choosing which function to use based on the type of the object
    for element in data["presentation"]:
        if (element["type"] == "title"):
            title_slide(element)
        elif (element["type"] == "text"):
            text_slide(element)
        elif (element["type"] == "list"):
            bullet_point_slide(element)
        elif (element["type"] == "picture"):
            picture_slide(element)
        elif (element["type"] == "plot"):
            plot_slide(element)

prs.save("test.pptx")